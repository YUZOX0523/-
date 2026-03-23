from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# カラーパレット
COLOR_BG = RGBColor(0xFF, 0xF8, 0xF0)       # クリーム
COLOR_PRIMARY = RGBColor(0xE8, 0x6A, 0x2E)  # オレンジ
COLOR_SECONDARY = RGBColor(0x5B, 0x8C, 0xDB) # ブルー
COLOR_DARK = RGBColor(0x2D, 0x2D, 0x2D)     # ほぼ黒
COLOR_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)     # 白
COLOR_ACCENT = RGBColor(0x27, 0xAE, 0x60)   # グリーン
COLOR_PINK = RGBColor(0xE9, 0x1E, 0x8C)     # ピンク
COLOR_YELLOW = RGBColor(0xFF, 0xC1, 0x07)   # 黄色

def set_bg(slide, color=COLOR_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, size=18, bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(0xFF, 0xF0, 0xE8))

    # 背景デコレーション
    add_rect(slide, 0, 0, 13.33, 1.5, COLOR_PRIMARY)
    add_rect(slide, 0, 6.0, 13.33, 1.5, COLOR_PRIMARY)

    # アクセント丸
    for x, y, size_i, col in [
        (11.5, 1.8, 1.2, RGBColor(0xFF, 0xD0, 0xA0)),
        (0.5, 5.0, 0.8, RGBColor(0xFF, 0xD0, 0xA0)),
        (12.0, 4.5, 0.6, COLOR_SECONDARY),
    ]:
        shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size_i), Inches(size_i))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()

    # タイトル
    add_textbox(slide, title, 1.0, 1.8, 11.0, 2.0, size=44, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    # サブタイトル
    add_textbox(slide, subtitle, 1.0, 3.8, 11.0, 1.0, size=22, color=COLOR_DARK, align=PP_ALIGN.CENTER, italic=True)
    # 日付
    add_textbox(slide, "2026年3月23日　Ver 1.0", 1.0, 6.1, 11.0, 0.7, size=16, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    return slide

def add_section_slide(prs, number, title, color=COLOR_PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, color)
    add_rect(slide, 0, 0, 13.33, 7.5, color)

    # 大きい番号
    add_textbox(slide, f"{number}", 1.0, 1.0, 4.0, 3.5, size=140, bold=True,
                color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 4.5, 2.5, 8.0, 2.5, size=40, bold=True,
                color=COLOR_LIGHT, align=PP_ALIGN.LEFT)
    return slide

def add_content_slide(prs, title, bullets, accent_color=COLOR_PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # 上部バー
    add_rect(slide, 0, 0, 13.33, 1.3, accent_color)
    # タイトル
    add_textbox(slide, title, 0.3, 0.1, 12.5, 1.1, size=28, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.LEFT)

    # 箇条書きエリア
    y = 1.6
    for i, (emoji, text) in enumerate(bullets):
        bg_col = RGBColor(0xFF,0xFF,0xFF) if i % 2 == 0 else RGBColor(0xF8, 0xF0, 0xE8)
        add_rect(slide, 0.4, y, 12.4, 0.65, bg_col)
        add_textbox(slide, emoji, 0.5, y+0.05, 0.6, 0.55, size=20)
        add_textbox(slide, text, 1.2, y+0.05, 11.2, 0.55, size=18, color=COLOR_DARK)
        y += 0.72

    return slide

def add_two_col_slide(prs, title, left_items, right_items, accent_color=COLOR_SECONDARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, 13.33, 1.3, accent_color)
    add_textbox(slide, title, 0.3, 0.1, 12.5, 1.1, size=28, bold=True, color=COLOR_LIGHT)

    # 左列ヘッダー
    add_rect(slide, 0.4, 1.5, 5.8, 0.55, accent_color)
    add_textbox(slide, left_items[0], 0.5, 1.55, 5.5, 0.45, size=18, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

    y = 2.1
    for item in left_items[1:]:
        add_rect(slide, 0.4, y, 5.8, 0.6, RGBColor(0xFF,0xFF,0xFF))
        add_textbox(slide, item, 0.5, y+0.05, 5.6, 0.5, size=17, color=COLOR_DARK)
        y += 0.65

    # 右列ヘッダー
    add_rect(slide, 7.0, 1.5, 5.8, 0.55, COLOR_PRIMARY)
    add_textbox(slide, right_items[0], 7.1, 1.55, 5.5, 0.45, size=18, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

    y = 2.1
    for item in right_items[1:]:
        add_rect(slide, 7.0, y, 5.8, 0.6, RGBColor(0xFF,0xFF,0xFF))
        add_textbox(slide, item, 7.1, y+0.05, 5.6, 0.5, size=17, color=COLOR_DARK)
        y += 0.65

    return slide

def add_step_slide(prs, title, steps, accent_color=COLOR_ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, 13.33, 1.3, accent_color)
    add_textbox(slide, title, 0.3, 0.1, 12.5, 1.1, size=28, bold=True, color=COLOR_LIGHT)

    x = 0.4
    step_w = (13.33 - 0.8) / len(steps)
    colors = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT, RGBColor(0x9B, 0x59, 0xB6), COLOR_PINK]

    for i, (num, step_title, detail) in enumerate(steps):
        col = colors[i % len(colors)]
        # ステップボックス
        add_rect(slide, x, 1.5, step_w - 0.15, 1.0, col)
        add_textbox(slide, f"STEP {num}", x+0.1, 1.55, step_w-0.3, 0.4, size=14, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(slide, step_title, x+0.1, 1.95, step_w-0.3, 0.5, size=16, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

        # 矢印（最後以外）
        if i < len(steps) - 1:
            add_textbox(slide, "▶", x + step_w - 0.2, 1.75, 0.4, 0.5, size=20, color=col, align=PP_ALIGN.CENTER)

        # 詳細テキスト
        add_rect(slide, x, 2.7, step_w-0.15, 4.4, RGBColor(0xFF,0xFF,0xFF))
        txBox = slide.shapes.add_textbox(Inches(x+0.1), Inches(2.8), Inches(step_w-0.35), Inches(4.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, line in enumerate(detail):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = COLOR_DARK
        x += step_w

    return slide

def add_qa_slide(prs, title, qas, accent_color=COLOR_SECONDARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, 13.33, 1.3, accent_color)
    add_textbox(slide, title, 0.3, 0.1, 12.5, 1.1, size=28, bold=True, color=COLOR_LIGHT)

    y = 1.5
    for q, a in qas:
        # Q
        add_rect(slide, 0.4, y, 0.55, 0.55, COLOR_PRIMARY)
        add_textbox(slide, "Q", 0.4, y, 0.55, 0.55, size=20, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(slide, q, 1.05, y, 11.8, 0.55, size=17, bold=True, color=COLOR_PRIMARY)
        y += 0.6
        # A
        add_rect(slide, 0.4, y, 0.55, 0.55, COLOR_ACCENT)
        add_textbox(slide, "A", 0.4, y, 0.55, 0.55, size=20, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(slide, a, 1.05, y, 11.8, 0.55, size=16, color=COLOR_DARK)
        y += 0.75

    return slide

def add_error_slide(prs, title, errors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, 13.33, 1.3, RGBColor(0xC0, 0x39, 0x2B))
    add_textbox(slide, title, 0.3, 0.1, 12.5, 1.1, size=28, bold=True, color=COLOR_LIGHT)

    y = 1.5
    for err_title, err_mean, err_fix in errors:
        add_rect(slide, 0.4, y, 12.4, 0.45, RGBColor(0xC0, 0x39, 0x2B))
        add_textbox(slide, f"🔴 {err_title}", 0.5, y+0.03, 12.0, 0.4, size=17, bold=True, color=COLOR_LIGHT)
        y += 0.48
        add_rect(slide, 0.4, y, 12.4, 0.45, RGBColor(0xFF, 0xEB, 0xEB))
        add_textbox(slide, f"意味: {err_mean}", 0.5, y+0.03, 12.0, 0.4, size=15, color=COLOR_DARK)
        y += 0.48
        add_rect(slide, 0.4, y, 12.4, 0.45, RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, f"対処: {err_fix}", 0.5, y+0.03, 12.0, 0.4, size=15, color=COLOR_DARK)
        y += 0.6

    return slide

def add_rules_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(0xFF, 0xF0, 0xE8))
    add_rect(slide, 0, 0, 13.33, 1.3, COLOR_YELLOW)
    add_textbox(slide, "🌟 上手に使うための黄金ルール", 0.3, 0.1, 12.5, 1.1, size=28, bold=True, color=COLOR_DARK)

    rules = [
        ("1", "具体的に伝える", "「文章書いて」より「○○宛ての○○についてのメールを書いて」"),
        ("2", "一度に一つずつ", "複数の作業は一つずつお願いする"),
        ("3", "気に入らなければやり直しを", "「もう一度、今度は〜で書いてください」"),
        ("4", "大切な情報は手動で保存", "回答はコピーしてメモ帳などに貼り付けて保存"),
        ("5", "個人情報は入力しない", "本名・住所・クレジットカード番号は絶対NG！"),
    ]

    colors_r = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT, RGBColor(0x9B, 0x59, 0xB6), COLOR_PINK]
    y = 1.5
    for num, rule, detail in rules:
        col = colors_r[int(num)-1]
        add_rect(slide, 0.4, y, 0.6, 0.65, col)
        add_textbox(slide, num, 0.4, y, 0.6, 0.65, size=24, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.1, y, 3.5, 0.65, col)
        add_textbox(slide, rule, 1.15, y+0.07, 3.4, 0.55, size=18, bold=True, color=COLOR_LIGHT)
        add_rect(slide, 4.7, y, 8.2, 0.65, RGBColor(0xFF,0xFF,0xFF))
        add_textbox(slide, detail, 4.8, y+0.07, 8.0, 0.55, size=16, color=COLOR_DARK)
        y += 0.75

    return slide

def add_ending_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_PRIMARY)
    add_rect(slide, 0, 0, 13.33, 7.5, COLOR_PRIMARY)

    add_textbox(slide, "🎉 さあ、始めましょう！", 0.5, 1.0, 12.0, 1.5, size=44, bold=True,
                color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "ClaudeCodeを起動して「こんにちは！」と話しかけるだけ", 0.5, 2.7, 12.0, 1.0,
                size=24, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

    steps = ["① ClaudeCodeを起動", "② 「こんにちは！」と入力", "③ それだけ！"]
    x = 1.5
    for s in steps:
        add_rect(slide, x, 4.0, 3.2, 0.9, RGBColor(0xFF,0xFF,0xFF))
        add_textbox(slide, s, x+0.1, 4.1, 3.0, 0.7, size=18, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
        x += 3.5

    add_textbox(slide, "このマニュアルは随時更新されます。困ったことがあればいつでも相談してください。",
                0.5, 5.2, 12.0, 0.7, size=15, color=RGBColor(0xFF, 0xD0, 0xB0), align=PP_ALIGN.CENTER)
    return slide

# ============================================================
# スライド作成
# ============================================================

# 1. タイトル
add_title_slide(prs,
    "🌸 ClaudeCode かんたんスタートガイド",
    "～IT音痴でも大丈夫！サルでもわかるマニュアル～")

# 2. 目次
add_content_slide(prs, "📋 目次", [
    ("1️⃣", "ClaudeCodeって何？"),
    ("2️⃣", "準備するもの（必要なもの）"),
    ("3️⃣", "インストール（入れ方）"),
    ("4️⃣", "はじめての起動"),
    ("5️⃣", "基本の使い方"),
    ("6️⃣", "実務での活用例"),
    ("7️⃣", "困ったときのQ&A"),
    ("8️⃣", "よくあるエラーと対処法"),
], COLOR_DARK)

# 3. セクション1
add_section_slide(prs, "01", "ClaudeCodeって何？", COLOR_PRIMARY)

# 4. ClaudeCode説明
add_two_col_slide(prs, "🤔 ClaudeCodeは「超賢いアシスタントさん」です！",
    ["普通のアシスタント", "☕ コーヒーを持ってくる", "📁 書類を整理する", "📅 スケジュール管理"],
    ["ClaudeCode", "✍️ 文章を書いてくれる", "📊 データを整理してくれる", "💻 プログラムを作ってくれる"],
    COLOR_PRIMARY)

# 5. 例え話
add_content_slide(prs, "💡 料理で例えると...", [
    ("👨‍🍳", "ClaudeCodeは「レシピを教えてくれる料理の先生」"),
    ("🍳", "「代わりに料理を作ってくれる料理人」でもあります"),
    ("😊", "あなたは「何を食べたいか」を言うだけでOK！"),
    ("🗣️", "日本語で普通に話しかければ理解してくれます"),
    ("⚡", "何十枚もの文書を数秒で処理できる高速アシスタント"),
], COLOR_PRIMARY)

# 6. セクション2
add_section_slide(prs, "02", "準備するもの", COLOR_SECONDARY)

# 7. チェックリスト
add_content_slide(prs, "✅ 準備するものチェックリスト", [
    ("💻", "パソコン（WindowsでもMacでもOK）"),
    ("🌐", "インターネット接続"),
    ("📧", "メールアドレス（Gmailで大丈夫）"),
    ("💳", "クレジットカードまたはデビットカード（月額費用の支払い用）"),
    ("💰", "無料プラン：毎月少しだけ使える（まずはこれでOK！）"),
    ("⭐", "有料プラン：月 $20（約3,000円）でたくさん使える"),
], COLOR_SECONDARY)

# 8. セクション3
add_section_slide(prs, "03", "インストール（入れ方）", COLOR_ACCENT)

# 9. インストール手順
add_step_slide(prs, "🎯 インストール5ステップ", [
    ("1", "アカウント作成", ["claude.ai にアクセス", "「Sign up」をクリック", "メールアドレスで登録"]),
    ("2", "Node.js を入れる", ["nodejs.org にアクセス", "緑のボタンをクリック", "「次へ次へ」で完了"]),
    ("3", "ターミナルを開く", ["Win: Windowsキー+R", "「cmd」と入力", "黒い画面が出たらOK！"]),
    ("4", "ClaudeCodeを入れる", ["ターミナルで以下を実行:", "npm install -g", "@anthropic-ai/claude-code"]),
    ("5", "APIキー設定", ["console.anthropic.com", "「API Keys」→作成", "sk-ant-... をコピー"]),
], COLOR_ACCENT)

# 10. ターミナルって何？
add_content_slide(prs, "🖥️ ターミナルって何？（怖くないよ！）", [
    ("💬", "パソコンへの「話しかけ窓口」です"),
    ("🔮", "コマンド（呪文）を入力してパソコンに命令する場所"),
    ("🪟", "Windowsの場合：「Windowsキー + R」→「cmd」と入力"),
    ("🍎", "Macの場合：「Command + スペース」→「ターミナル」と入力"),
    ("⚠️", "APIキーは他人に絶対見せない！（銀行の暗証番号と同じ）"),
], COLOR_ACCENT)

# 11. セクション4
add_section_slide(prs, "04", "はじめての起動", RGBColor(0x9B, 0x59, 0xB6))

# 12. 起動方法
add_content_slide(prs, "🚀 ClaudeCodeを起動してみよう", [
    ("1️⃣", "ターミナルを開く"),
    ("2️⃣", "「claude」と入力してEnterキーを押す"),
    ("3️⃣", "Welcome to Claude Code! と表示されたら成功！🎉"),
    ("4️⃣", "「こんにちは！」と入力してEnterを押してみよう"),
    ("5️⃣", "ClaudeCodeが返事をしてくれます！"),
], RGBColor(0x9B, 0x59, 0xB6))

# 13. セクション5
add_section_slide(prs, "05", "基本の使い方", COLOR_PINK)

# 14. 話しかけ方
add_content_slide(prs, "💬 話しかけ方のコツ", [
    ("❌", "NG：「文章書いて」（曖昧すぎる）"),
    ("✅", "OK：「取引先への謝罪メールを書いてください。商品の納期が2日遅れました。丁寧な文体で」"),
    ("👤", "「誰に」を伝える"),
    ("📝", "「何を」を伝える"),
    ("🎨", "「どんな感じで」を伝える"),
    ("💡", "この3つを伝えるとぴったりの答えが返ってきます！"),
], COLOR_PINK)

# 15. セクション6
add_section_slide(prs, "06", "実務での活用例", COLOR_SECONDARY)

# 16. 活用例1
add_content_slide(prs, "📧 実務活用例（前半）", [
    ("📧", "例1：ビジネスメールを書いてもらう → 相手・内容・日時を伝えるだけ"),
    ("📊", "例2：Excelのデータを整理してもらう → データをコピー＆ペーストするだけ"),
    ("✍️", "例3：文章の校正をしてもらう → 文章を貼り付けて「校正して」と言うだけ"),
    ("📋", "例4：議事録を作ってもらう → メモを貼り付けて「議事録にして」と言うだけ"),
], COLOR_SECONDARY)

# 17. 活用例2
add_content_slide(prs, "🔍 実務活用例（後半）", [
    ("🔍", "例5：調べ物 → 「中学生でもわかるように説明して」"),
    ("💡", "例6：アイデア出し → 「SNS投稿のアイデアを10個出して」"),
    ("📈", "例7：プログラム作成 → 「sample.csvの売上合計を計算するプログラムを」"),
    ("🎯", "ポイント：具体的に伝えるほど、良い答えが返ってくる！"),
    ("🔄", "気に入らなければ「もう一度、今度は〜で」と言えばOK！"),
], COLOR_SECONDARY)

# 18. セクション7
add_section_slide(prs, "07", "困ったときのQ&A", RGBColor(0x16, 0xA0, 0x85))

# 19. Q&A
add_qa_slide(prs, "❓ よくある質問 Q&A",
    [
        ("日本語で話しかけても大丈夫？", "大丈夫！ClaudeCodeは日本語をちゃんと理解します"),
        ("答えが長すぎて読めない…", "「もっと短くまとめて。3行でお願いします」と言えばOK！"),
        ("答えが難しくてわからない…", "「もっと簡単に。小学生でもわかるように」と言えばOK！"),
        ("間違った答えが返ってきた…", "「それは違います。○○については××が正しいです」と訂正して！"),
        ("作業を止めたい！", "Ctrl + C を同時に押すと止まります"),
    ], RGBColor(0x16, 0xA0, 0x85))

# 20. セクション8
add_section_slide(prs, "08", "よくあるエラーと対処法", RGBColor(0xC0, 0x39, 0x2B))

# 21. エラー対処
add_error_slide(prs, "🔴 よくあるエラーと対処法",
    [
        ("「command not found」", "ClaudeCodeが見つかりません",
         "ターミナルを再起動 → npm install -g @anthropic-ai/claude-code を再実行"),
        ("「API key not found」", "APIキーが設定されていません",
         "set ANTHROPIC_API_KEY=sk-ant-xxx（Win）/ export ANTHROPIC_API_KEY=sk-ant-xxx（Mac）"),
        ("文字化けして読めない", "文字の表示がおかしい",
         "「日本語で回答してください」と入力。直らない場合はターミナルを再起動"),
        ("動きが遅い・固まった", "処理中かインターネット不安定",
         "1〜2分待つ → Ctrl+C で止めてやり直す"),
    ])

# 22. 黄金ルール
add_rules_slide(prs)

# 23. まとめ・終わり
add_ending_slide(prs)

# 保存
output_path = "/home/user/-/ClaudeCode完全マニュアル_超初心者向け.pptx"
prs.save(output_path)
print(f"✅ 完成！→ {output_path}")
print(f"   スライド数: {len(prs.slides)} 枚")
